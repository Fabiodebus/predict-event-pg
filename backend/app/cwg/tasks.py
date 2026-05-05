"""CWG (Customer Workspace & GTM Thesis) Celery tasks.

Auto-discovered by app.celery_app via celery_app.autodiscover_tasks(["app.cwg"]).
"""

from __future__ import annotations

import asyncio
import logging
from functools import lru_cache
from io import BytesIO
from typing import Any
from uuid import UUID

import boto3
from celery import Task
from pydantic import BaseModel, Field, ValidationError

from app.celery_app import celery_app
from app.config import settings
from app.integrations.crustdata import CrustdataClient
from app.llm.anthropic_provider import AnthropicProvider
from app.llm.openai_provider import OpenAIProvider
from app.llm.provider import LLMMessage
from app.llm.roles import LLMRole
from app.llm.router import LLMRouter
from app.models.best_customer import BestCustomer, EnrichmentStatus
from app.models.sales_material import (
    ContentType,
    ExtractionStatus,
    SalesMarketingMaterial,
)
from app.tasks.base import IdempotentTask, _open_session

logger = logging.getLogger(__name__)


# Task-level retries are disabled because BaseIntegrationClient already retries
# transient 5xx via tenacity, and the BC.enrichment_status field is the source
# of truth for business-level success/failure (orthogonal to Job lifecycle).
@celery_app.task(
    bind=True,
    base=IdempotentTask,
    name="app.cwg.enrich_best_customer",
    autoretry_for=(),
    max_retries=0,
)
def enrich_best_customer(self: Task, *, best_customer_id: str) -> dict[str, Any]:
    """Enrich a single BestCustomer entry via Crustdata.

    Reads the BestCustomer, picks the best identifier (LinkedIn URL > domain >
    company_name), calls Crustdata's screener, and writes industry /
    employee_count / hq_country back. Sets enrichment_status to enriched on
    success, failed on any of: no identifier present, no Crustdata match,
    Crustdata error after retries.
    """
    return asyncio.run(_run_enrich_best_customer(UUID(best_customer_id)))


def _select_identifier(bc: BestCustomer) -> tuple[str, str] | None:
    """Identifier preference: LinkedIn URL > domain > company_name.

    Returns (kwarg_name_for_screener_company, identifier_value) or None.
    """
    if bc.linkedin_url:
        return ("company_linkedin_urls", bc.linkedin_url)
    if bc.domain:
        return ("company_domains", bc.domain)
    if bc.company_name:
        return ("company_names", bc.company_name)
    return None


def _extract_industry(profile: dict[str, Any]) -> str | None:
    industries = profile.get("linkedin_industries") or []
    if isinstance(industries, list) and industries:
        return industries[0]
    # Fallback: some Crustdata responses use the `industry` field directly.
    industry = profile.get("industry")
    return industry if isinstance(industry, str) else None


def _extract_employee_count(profile: dict[str, Any]) -> int | None:
    headcount = profile.get("headcount")
    if isinstance(headcount, dict):
        latest = headcount.get("latest_count")
        if isinstance(latest, int):
            return latest
    return None


async def _run_enrich_best_customer(best_customer_id: UUID) -> dict[str, Any]:
    async with _open_session() as session:
        bc = await session.get(BestCustomer, best_customer_id)
        if bc is None:
            return {"status": "skipped", "reason": "best_customer_not_found"}

        identifier = _select_identifier(bc)
        if identifier is None:
            bc.enrichment_status = EnrichmentStatus.FAILED
            await session.commit()
            return {"status": "failed", "reason": "no_identifier"}

        client = CrustdataClient(api_key=settings.crustdata_api_key)
        kwarg_name, value = identifier
        try:
            profiles = await client.screener_company(**{kwarg_name: [value]})
        except Exception as exc:
            bc.enrichment_status = EnrichmentStatus.FAILED
            await session.commit()
            logger.exception(
                "crustdata enrichment failed for best_customer_id=%s",
                best_customer_id,
            )
            return {"status": "failed", "reason": f"crustdata_error: {type(exc).__name__}"}

        profile = profiles[0] if profiles else None
        if profile is None:
            bc.enrichment_status = EnrichmentStatus.FAILED
            await session.commit()
            return {"status": "failed", "reason": "no_match"}

        bc.industry = _extract_industry(profile)
        bc.employee_count = _extract_employee_count(profile)
        bc.hq_country = profile.get("hq_country")
        bc.enrichment_status = EnrichmentStatus.ENRICHED
        await session.commit()
        return {"status": "ok"}


# ---------------------------------------------------------------------------
# Sales / marketing material extraction
# ---------------------------------------------------------------------------


# Truncate the document text before sending to the LLM. ~24k chars maps to
# roughly 6k tokens for English; well within structured_extraction model
# context. Extraction quality saturates well before that for sales material.
_MAX_EXTRACTION_CHARS = 24_000

_EXTRACTION_SYSTEM_PROMPT = """You are extracting structured information from a B2B company's sales or marketing material.

Return ONLY a single valid JSON object with these keys (no commentary, no markdown fences):
- solution_description (string or null): What the company sells, in 1-3 sentences. Null if unclear.
- proof_points (array of strings): Quantified claims, customer wins, metrics, awards. Each item is one short string.
- use_cases (array of strings): Concrete customer scenarios where the solution applies.
- customer_references (array of strings): Named customers, case studies, or testimonials.
- communication_style_indicators (array of strings): Phrases or word choices that characterize the company's tone.

Use empty arrays for absent categories. Output must be valid JSON parseable directly with json.loads.
"""


class ExtractedContent(BaseModel):
    """Schema persisted to SalesMarketingMaterial.extracted_content."""

    solution_description: str | None = None
    proof_points: list[str] = Field(default_factory=list)
    use_cases: list[str] = Field(default_factory=list)
    customer_references: list[str] = Field(default_factory=list)
    communication_style_indicators: list[str] = Field(default_factory=list)


@celery_app.task(
    bind=True,
    base=IdempotentTask,
    name="app.cwg.extract_material",
    autoretry_for=(),
    max_retries=0,
)
def extract_material(self: Task, *, material_id: str) -> dict[str, Any]:
    """Extract structured content from a sales/marketing material.

    Steps: load text (S3 download + parse, or raw_text passthrough) →
    LLM extraction with the structured_extraction role → persist
    extracted_content + extraction_status.
    """
    return asyncio.run(_run_extract_material(UUID(material_id)))


async def _run_extract_material(material_id: UUID) -> dict[str, Any]:
    async with _open_session() as session:
        m = await session.get(SalesMarketingMaterial, material_id)
        if m is None:
            return {"status": "skipped", "reason": "material_not_found"}

        try:
            text = await _load_material_text(m)
        except Exception as exc:
            m.extraction_status = ExtractionStatus.FAILED
            await session.commit()
            logger.exception("material text load failed for material_id=%s", material_id)
            return {"status": "failed", "reason": f"load_error: {type(exc).__name__}"}

        if not text.strip():
            m.extraction_status = ExtractionStatus.FAILED
            await session.commit()
            return {"status": "failed", "reason": "empty_content"}

        try:
            extracted = await _extract_via_llm(text)
        except (ValidationError, ValueError) as exc:
            m.extraction_status = ExtractionStatus.FAILED
            await session.commit()
            logger.warning(
                "LLM returned non-conformant extraction for material_id=%s: %s",
                material_id,
                exc,
            )
            return {"status": "failed", "reason": "llm_output_invalid"}
        except Exception as exc:
            m.extraction_status = ExtractionStatus.FAILED
            await session.commit()
            logger.exception("LLM extraction failed for material_id=%s", material_id)
            return {"status": "failed", "reason": f"llm_error: {type(exc).__name__}"}

        m.extracted_content = extracted.model_dump()
        m.extraction_status = ExtractionStatus.EXTRACTED
        await session.commit()
        return {"status": "ok"}


async def _load_material_text(m: SalesMarketingMaterial) -> str:
    if m.content_type == ContentType.TEXT:
        return m.raw_text or ""
    if not m.s3_key:
        raise ValueError(f"non-text material {m.id} has no s3_key")
    body = await _download_s3_object(m.s3_key)
    if m.content_type == ContentType.PDF:
        return _parse_pdf(body)
    if m.content_type == ContentType.DOCX:
        return _parse_docx(body)
    if m.content_type == ContentType.PPTX:
        return _parse_pptx(body)
    raise ValueError(f"unsupported content_type: {m.content_type}")


async def _download_s3_object(key: str) -> bytes:
    """Download an S3 object body. Boto3 is sync, so run in default executor."""
    def _sync() -> bytes:
        client = boto3.client("s3", region_name=settings.aws_region)
        obj = client.get_object(Bucket=settings.aws_s3_bucket, Key=key)
        return obj["Body"].read()

    return await asyncio.get_running_loop().run_in_executor(None, _sync)


def _parse_pdf(body: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(body))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _parse_docx(body: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(body))
    return "\n".join(p.text for p in doc.paragraphs)


def _parse_pptx(body: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(body))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
    return "\n".join(parts)


@lru_cache(maxsize=1)
def _build_router() -> LLMRouter:
    """Module-level cached router so tasks reuse the underlying HTTP clients."""
    return LLMRouter(
        providers=[
            AnthropicProvider(api_key=settings.anthropic_api_key),
            OpenAIProvider(api_key=settings.openai_api_key),
        ],
        role_models=settings.llm_role_models,
    )


async def _extract_via_llm(text: str) -> ExtractedContent:
    truncated = text[:_MAX_EXTRACTION_CHARS]
    response = await _build_router().acomplete(
        role=LLMRole.STRUCTURED_EXTRACTION,
        messages=[
            LLMMessage(role="system", content=_EXTRACTION_SYSTEM_PROMPT),
            LLMMessage(role="user", content=truncated),
        ],
        max_tokens=2048,
        temperature=0.0,
    )
    return ExtractedContent.model_validate_json(response.text)
